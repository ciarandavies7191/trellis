"""generate_documents — assemble approved step outputs into DOCX/PDF/PPTX/XLSX.

No LLM calls — pure document assembly from structured markdown step outputs.

Optional dependencies (install via ``pip install trellis-pipelines[credit-dd]``):
    python-docx, python-pptx, openpyxl, weasyprint (or reportlab for PDF).

If a dependency is missing the corresponding format is skipped and the path key
is absent from the return dict; an error message is appended to ``errors``.

Registration name: ``generate_documents``
"""

from __future__ import annotations

import logging
import os
import re
from pathlib import Path
from typing import Any, Dict, List, Optional

from trellis.tools.base import BaseTool, ToolInput, ToolOutput

logger = logging.getLogger(__name__)

# Step-output assembly order (matches credit-DD memo structure)
_STEP_ORDER = [
    ("step_02_output", "Historical Financials"),
    ("step_03_output", "Obligor Overview"),
    ("step_04_output", "Industry & Peers"),
    ("step_05_output", "Capital Structure"),
    ("step_06_output", "Repayment & Liquidity"),
    ("step_07_output", "Pro Forma / Stress"),
    ("step_08_output", "Projections & DCF"),
    ("step_09_output", "Risks & Mitigants"),
    ("step_12_output", "Rating & Classification"),
    ("step_10_output", "Recommendation"),
]

# XLSX tab names mapped to step keys that supply their content
_XLSX_TABS = [
    ("Financial Summary",   "step_02_output"),
    ("Peer Comparison",     "step_04_output"),
    ("Capital Structure",   "step_05_output"),
    ("Pro Forma Stress",    "step_07_output"),
    ("DCF Model",           "step_08_output"),
]

# PPTX slide layout (title, source key)
_PPTX_SLIDES = [
    ("Cover",                        None),
    ("Obligor Profile",              "step_03_output"),
    ("Historical Financials",        "step_02_output"),
    ("Capital Structure & Maturity", "step_05_output"),
    ("Peer Comparison",              "step_04_output"),
    ("Repayment & Coverage",         "step_06_output"),
    ("Pro Forma + Stress",           "step_07_output"),
    ("Risk Heat Map",                "step_09_output"),
    ("Rating & OCC Classification",  "step_12_output"),
    ("Recommendation",               "step_10_output"),
]


def _md_to_plain(text: str, max_chars: int = 2000) -> str:
    """Strip markdown syntax to plain text for PPTX/DOCX body."""
    if not text:
        return ""
    # Remove code fences, headings markers, horizontal rules, bold/italic
    plain = re.sub(r"^#+\s*", "", text, flags=re.MULTILINE)
    plain = re.sub(r"\*{1,3}([^*]+)\*{1,3}", r"\1", plain)
    plain = re.sub(r"_{1,3}([^_]+)_{1,3}", r"\1", plain)
    plain = re.sub(r"^[-*_]{3,}$", "", plain, flags=re.MULTILINE)
    plain = re.sub(r"^```.*?```", "", plain, flags=re.DOTALL)
    plain = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", plain)
    plain = re.sub(r"\n{3,}", "\n\n", plain)
    return plain.strip()[:max_chars]


def _build_docx(doc: Any, step_outputs: dict, deal_metadata: dict) -> None:
    from docx.shared import Pt, RGBColor  # type: ignore[import]
    from docx.enum.text import WD_ALIGN_PARAGRAPH  # type: ignore[import]

    ticker = deal_metadata.get("ticker", "")
    borrower = deal_metadata.get("borrower_name", ticker)

    # Cover page
    cover = doc.add_heading(f"Credit Due-Diligence Memo", level=0)
    cover.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph(
        f"{borrower}  |  Facility: "
        f"${deal_metadata.get('facility_amount_mm', '')}MM  |  "
        f"{deal_metadata.get('facility_product', '')}",
    )
    doc.add_page_break()

    for step_key, section_title in _STEP_ORDER:
        content = step_outputs.get(step_key, "")
        if not content:
            continue
        doc.add_heading(section_title, level=1)
        for para in content.split("\n\n"):
            para = para.strip()
            if not para:
                continue
            # Render section headings
            if para.startswith("### "):
                doc.add_heading(para[4:], level=3)
            elif para.startswith("## "):
                doc.add_heading(para[3:], level=2)
            else:
                doc.add_paragraph(para)
        doc.add_page_break()


def _build_pptx(prs: Any, step_outputs: dict, deal_metadata: dict) -> None:
    from pptx.util import Inches, Pt  # type: ignore[import]
    from pptx.enum.text import PP_ALIGN  # type: ignore[import]

    ticker = deal_metadata.get("ticker", "")
    borrower = deal_metadata.get("borrower_name", ticker)

    blank_layout = prs.slide_layouts[6]  # blank
    title_layout = prs.slide_layouts[0]

    for slide_title, step_key in _PPTX_SLIDES:
        if step_key is None:
            # Cover slide
            slide = prs.slides.add_slide(title_layout)
            slide.shapes.title.text = f"Credit Due-Diligence\n{borrower}"
            if slide.placeholders and len(slide.placeholders) > 1:
                slide.placeholders[1].text = (
                    f"Facility: ${deal_metadata.get('facility_amount_mm','')}MM  "
                    f"{deal_metadata.get('facility_product','')}"
                )
            continue

        slide = prs.slides.add_slide(blank_layout)
        # Title text box
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
        tf = txBox.text_frame
        tf.text = slide_title
        tf.paragraphs[0].runs[0].font.size = Pt(20)
        tf.paragraphs[0].runs[0].font.bold = True

        # Body text box
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
        body.text_frame.word_wrap = True
        content = step_outputs.get(step_key, "")
        body.text_frame.text = _md_to_plain(content, max_chars=1200)


def _build_xlsx(wb: Any, step_outputs: dict, deal_metadata: dict) -> None:
    from openpyxl.styles import Font, PatternFill, Alignment  # type: ignore[import]
    from openpyxl.utils import get_column_letter  # type: ignore[import]

    header_font = Font(bold=True, color="FFFFFF")
    header_fill = PatternFill(fill_type="solid", fgColor="1F4E79")

    # Remove default blank sheet
    default = wb.active
    if default is not None:
        wb.remove(default)

    for tab_name, step_key in _XLSX_TABS:
        content = step_outputs.get(step_key, "")
        if not content:
            ws = wb.create_sheet(title=tab_name)
            ws["A1"] = "No data available"
            continue

        ws = wb.create_sheet(title=tab_name)
        row = 1

        # Parse markdown tables into rows
        in_table = False
        for line in content.split("\n"):
            line = line.strip()
            if line.startswith("|") and line.endswith("|"):
                # Skip separator rows (|---|---|)
                if re.match(r"^\|[-: |]+\|$", line):
                    continue
                cols = [c.strip() for c in line.strip("|").split("|")]
                for col_idx, cell_val in enumerate(cols, start=1):
                    cell = ws.cell(row=row, column=col_idx, value=cell_val)
                    if row == 1 or not in_table:
                        cell.font = header_font
                        cell.fill = header_fill
                    cell.alignment = Alignment(wrap_text=True)
                in_table = True
                row += 1
            else:
                if in_table:
                    row += 1  # blank row after table
                in_table = False
                if line.startswith("## ") or line.startswith("### "):
                    ws.cell(row=row, column=1, value=line.lstrip("#").strip()).font = Font(bold=True)
                    row += 1

        # Auto-size columns (rough estimate)
        for col in ws.columns:
            max_len = max((len(str(c.value or "")) for c in col), default=10)
            ws.column_dimensions[get_column_letter(col[0].column)].width = min(max_len + 2, 60)


class GenerateDocumentsTool(BaseTool):
    """Assemble all approved step outputs into DOCX, PDF, PPTX, and XLSX files.

    Requires optional dependencies (``pip install trellis-pipelines[credit-dd]``).
    Missing dependencies cause the corresponding format to be skipped; the
    ``errors`` list in the response describes what was skipped and why.

    Registration name: ``generate_documents``
    """

    def __init__(self, name: str = "generate_documents") -> None:
        super().__init__(
            name,
            "Assemble step outputs into DOCX, PDF, PPTX, XLSX deliverables",
        )

    def execute(
        self,
        step_outputs: dict,
        deal_metadata: dict,
        output_dir: str = "./outputs",
        formats: Optional[List[str]] = None,
        **_: Any,
    ) -> Dict[str, Any]:
        """
        Args:
            step_outputs:  Keys step_02_output … step_12_output (markdown strings).
            deal_metadata: ticker, borrower_name, facility_amount_mm, etc.
            output_dir:    Directory path for output files (created if absent).
            formats:       Which formats to generate: subset of [docx, pdf, pptx, xlsx].
                           Defaults to all four.

        Returns:
            dict with optional keys docx, pdf, pptx, xlsx (absolute paths) and
            an ``errors`` list for any formats that failed or were skipped.
        """
        formats = [f.lower() for f in (formats or ["docx", "pdf", "pptx", "xlsx"])]
        out_dir = Path(output_dir)
        out_dir.mkdir(parents=True, exist_ok=True)
        ticker = deal_metadata.get("ticker", "UNKNOWN")
        result: Dict[str, Any] = {}
        errors: List[str] = []

        # ── DOCX ──────────────────────────────────────────────────────────
        if "docx" in formats:
            try:
                from docx import Document  # type: ignore[import]
                doc = Document()
                _build_docx(doc, step_outputs, deal_metadata)
                docx_path = str(out_dir / f"{ticker}_credit_memo.docx")
                doc.save(docx_path)
                result["docx"] = docx_path
                logger.info("generate_documents: saved %s", docx_path)
            except ImportError:
                errors.append(
                    "docx: python-docx not installed. "
                    "Run: pip install trellis-pipelines[credit-dd]"
                )
            except Exception as exc:
                errors.append(f"docx: {exc}")

        # ── PDF ───────────────────────────────────────────────────────────
        if "pdf" in formats:
            try:
                # Prefer weasyprint; fall back to a plain-text PDF via reportlab
                pdf_path = str(out_dir / f"{ticker}_credit_memo.pdf")
                combined_md = "\n\n".join(
                    f"## {title}\n{step_outputs.get(key, '')}"
                    for key, title in _STEP_ORDER
                    if step_outputs.get(key)
                )
                try:
                    import weasyprint  # type: ignore[import]
                    html_content = (
                        "<html><body style='font-family:Arial;margin:40px'>"
                        + re.sub(
                            r"^(#{1,3})\s*(.+)$",
                            lambda m: f"<h{len(m.group(1))}>{m.group(2)}</h{len(m.group(1))}>",
                            combined_md,
                            flags=re.MULTILINE,
                        ).replace("\n\n", "<br><br>")
                        + "</body></html>"
                    )
                    weasyprint.HTML(string=html_content).write_pdf(pdf_path)
                except ImportError:
                    from reportlab.lib.pagesizes import letter  # type: ignore[import]
                    from reportlab.platypus import SimpleDocTemplate, Paragraph  # type: ignore[import]
                    from reportlab.lib.styles import getSampleStyleSheet  # type: ignore[import]

                    doc_pdf = SimpleDocTemplate(pdf_path, pagesize=letter)
                    styles = getSampleStyleSheet()
                    story = []
                    for key, title in _STEP_ORDER:
                        content = step_outputs.get(key, "")
                        if not content:
                            continue
                        story.append(Paragraph(title, styles["Heading1"]))
                        for para in _md_to_plain(content, max_chars=50000).split("\n\n"):
                            if para.strip():
                                story.append(Paragraph(para.strip(), styles["Normal"]))
                    doc_pdf.build(story)

                result["pdf"] = pdf_path
                logger.info("generate_documents: saved %s", pdf_path)
            except Exception as exc:
                errors.append(f"pdf: {exc}")

        # ── PPTX ──────────────────────────────────────────────────────────
        if "pptx" in formats:
            try:
                from pptx import Presentation  # type: ignore[import]
                prs = Presentation()
                _build_pptx(prs, step_outputs, deal_metadata)
                pptx_path = str(out_dir / f"{ticker}_credit_summary.pptx")
                prs.save(pptx_path)
                result["pptx"] = pptx_path
                logger.info("generate_documents: saved %s", pptx_path)
            except ImportError:
                errors.append(
                    "pptx: python-pptx not installed. "
                    "Run: pip install trellis-pipelines[credit-dd]"
                )
            except Exception as exc:
                errors.append(f"pptx: {exc}")

        # ── XLSX ──────────────────────────────────────────────────────────
        if "xlsx" in formats:
            try:
                import openpyxl  # type: ignore[import]
                wb = openpyxl.Workbook()
                _build_xlsx(wb, step_outputs, deal_metadata)
                xlsx_path = str(out_dir / f"{ticker}_data_tables.xlsx")
                wb.save(xlsx_path)
                result["xlsx"] = xlsx_path
                logger.info("generate_documents: saved %s", xlsx_path)
            except ImportError:
                errors.append(
                    "xlsx: openpyxl not installed. "
                    "Run: pip install trellis-pipelines[credit-dd]"
                )
            except Exception as exc:
                errors.append(f"xlsx: {exc}")

        result["errors"] = errors
        return result

    def get_inputs(self) -> Dict[str, ToolInput]:
        return {
            "step_outputs": ToolInput(
                name="step_outputs",
                description="Keys step_02_output … step_12_output (markdown strings)",
                required=True,
            ),
            "deal_metadata": ToolInput(
                name="deal_metadata",
                description="ticker, borrower_name, facility_amount_mm, facility_product, etc.",
                required=True,
            ),
            "output_dir": ToolInput(
                name="output_dir",
                description="Directory for output files (created if absent)",
                required=False,
                default="./outputs",
            ),
            "formats": ToolInput(
                name="formats",
                description="Subset of [docx, pdf, pptx, xlsx] to generate (default: all)",
                required=False,
                default=None,
            ),
        }

    def get_output(self) -> ToolOutput:
        return ToolOutput(
            name="files",
            description="Paths to generated files (docx, pdf, pptx, xlsx) + errors list",
            type_="object",
        )
